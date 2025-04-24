import streamlit as st
st.set_page_config(page_title="Document Q&A (Claude 3)", page_icon="📘")

from htmltemplates import css, render_user_message, render_bot_message

import os
import io
import re
import json
import requests
import pandas as pd
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import anthropic

# Load environment variables
load_dotenv()
ANTHROPIC_API_KEY = os.getenv("ANTHROPIC_API_KEY")

# Tesseract path
pytesseract.pytesseract.tesseract_cmd = r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe"

# Claude client
client = anthropic.Anthropic(api_key=ANTHROPIC_API_KEY)

# Claude call with chunking + search
def chunk_text(text, chunk_size=3000, overlap=500):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = words[i:i+chunk_size]
        chunks.append(" ".join(chunk))
        i += chunk_size - overlap
    return chunks

def find_relevant_chunk(chunks, question):
    question_keywords = set(question.lower().split())
    scores = []
    for chunk in chunks:
        chunk_words = set(chunk.lower().split())
        score = len(question_keywords & chunk_words)
        scores.append(score)
    best_index = scores.index(max(scores))
    return chunks[best_index]

def call_claude(question, full_context):
    chunks = chunk_text(full_context)
    relevant_chunk = find_relevant_chunk(chunks, question)
    response = client.messages.create(
        model="claude-3-haiku-20240307",
        max_tokens=1024,
        temperature=0.2,
        messages=[
            {"role": "user", "content": f"{relevant_chunk}\n\nQuestion: {question}"}
        ]
    )
    return response.content[0].text.strip()

# Extract text from various file types
def extract_text_from_file(file):
    ext = file.name.split(".")[-1].lower()
    text = ""
    try:
        if ext == "pdf":
            pdf_reader = PdfReader(file)
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
        elif ext == "docx":
            doc = Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == "pptx":
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif ext == "xlsx":
            df = pd.read_excel(file)
            text += df.to_string(index=False)
        elif ext == "csv":
            file.seek(0)
            df = pd.read_csv(file)
            text += df.to_string(index=False)
        elif ext == "json":
            file.seek(0)
            data = json.load(file)
            text += json.dumps(data, indent=2)
        elif ext == "txt":
            stringio = io.StringIO(file.getvalue().decode("utf-8"))
            text += stringio.read()
        elif ext in ["png", "jpg", "jpeg"]:
            image = Image.open(file)
            text += pytesseract.image_to_string(image)

        urls = re.findall(r'https?://\S+', text)
        for url in urls:
            try:
                response = requests.get(url, timeout=5)
                soup = BeautifulSoup(response.text, 'html.parser')
                paragraphs = soup.find_all('p')
                external_text = "\n".join(p.text for p in paragraphs)
                text += f"\n--- Fetched from {url} ---\n{external_text}\n"
            except:
                text += f"\n[Failed to fetch {url}]\n"
    except Exception as e:
        text += f"\n[Error reading {file.name}: {e}]"
    return text

def get_combined_text(files):
    return "\n".join(
        f"\n--- {file.name} ---\n{text}\n"
        for file in files
        if (text := extract_text_from_file(file)).strip() and "Error reading" not in text
    )

def main():
    st.markdown(css, unsafe_allow_html=True)
    st.title("📄 Chat with Your Documents (Claude 3 Haiku)")

    if "doc_context" not in st.session_state:
        st.session_state.doc_context = ""
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    for role, message in st.session_state.chat_history:
        bubble = render_user_message(message) if role == "user" else render_bot_message(message)
        st.markdown(bubble, unsafe_allow_html=True)

    if prompt := st.chat_input("Ask a question about your uploaded document(s):"):
        if not st.session_state.doc_context:
            st.warning("Please upload and process a document first.")
        else:
            st.chat_message("user").markdown(prompt)
            with st.spinner("Thinking..."):
                try:
                    answer = call_claude(prompt, st.session_state.doc_context)
                    st.chat_message("assistant").markdown(answer)
                    st.session_state.chat_history.append(("user", prompt))
                    st.session_state.chat_history.append(("assistant", answer))
                except Exception as e:
                    st.error(f"Error generating answer: {str(e)}")

    # Sidebar
    st.sidebar.subheader("Upload Files")
    uploaded_files = st.sidebar.file_uploader(
        "Supported: PDF, DOCX, PPTX, XLSX, PNG, JPG, CSV, JSON, TXT",
        type=["pdf", "docx", "pptx", "xlsx", "png", "jpg", "jpeg", "csv", "json", "txt"],
        accept_multiple_files=True
    )

    if st.sidebar.button("Process"):
        if uploaded_files:
            with st.spinner("Extracting text..."):
                raw_text = get_combined_text(uploaded_files)
                if not raw_text.strip():
                    st.error("No usable text found.")
                else:
                    st.session_state.doc_context = raw_text
                    st.session_state.chat_history = []
                    st.success("Document loaded successfully!")
        else:
            st.warning("Please upload at least one file.")

    if st.sidebar.button("Clear"):
        st.session_state.doc_context = ""
        st.session_state.chat_history = []
        st.success("Cleared uploaded document and chat.")

if __name__ == '__main__':
    main()
